from PyPDF2 import PdfReader
import os
from langchain.text_splitter import (
    RecursiveCharacterTextSplitter,
    CharacterTextSplitter,
    TokenTextSplitter
)
from typing import List, Dict, Literal, Union
import pptx
from docx import Document
import pandas as pd
import zipfile
from io import BytesIO, StringIO

ChunkingStrategy = Literal["recursive", "character", "token"]

def load_file_content(file_path: str) -> Union[str, List[Dict]]:
    """Load content from various file formats and return extracted text or structured data.
    
    Args:
        file_path: Path to the file to be loaded.
        
    Returns:
        For PDF: Returns list of dicts with text and page numbers
        For other text-based files (PPTX, DOCX, CSV): Returns extracted text.
        For Excel: Returns list of dictionaries representing each sheet's data.
        For ZIP: Returns list of dictionaries with extracted content from supported files within.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    if ext == '.pdf':
        return load_pdf_content_with_pages(file_path)
    elif ext == '.pptx':
        return load_pptx_content(file_path)
    elif ext == '.docx':
        return load_docx_content(file_path)
    elif ext in ('.xlsx', '.xls'):
        return load_excel_content(file_path)
    elif ext == '.csv':
        return load_csv_content(file_path)
    elif ext == '.zip':
        return load_zip_content(file_path)
    else:
        raise ValueError(f"Unsupported file format: {ext}")

def load_pdf_content_with_pages(file_path: str) -> List[Dict]:
    """Extract text from PDF file with page numbers."""
    reader = PdfReader(file_path)
    pages_content = []
    for page_num, page in enumerate(reader.pages, start=1):
        text = page.extract_text()
        if text.strip():
            pages_content.append({
                "text": text,
                "page": page_num
            })
    return pages_content

def load_pptx_content(file_path: str) -> str:
    """Extract text from PowerPoint file."""
    presentation = pptx.Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text).strip()

def load_docx_content(file_path: str) -> str:
    """Extract text from Word document."""
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs]).strip()

def load_excel_content(file_path: str) -> List[Dict]:
    """Extract data from Excel file (all sheets)."""
    xls = pd.ExcelFile(file_path)
    sheets_data = []
    for sheet_name in xls.sheet_names:
        df = pd.read_excel(xls, sheet_name=sheet_name)
        sheets_data.append({
            "sheet_name": sheet_name,
            "data": df.to_dict(orient="records")
        })
    return sheets_data

def load_csv_content(file_path: str) -> str:
    """Extract content from CSV file."""
    df = pd.read_csv(file_path)
    return df.to_string(index=False)

def load_zip_content(file_path: str) -> List[Dict]:
    """Extract content from supported files within a ZIP archive."""
    extracted_contents = []
    with zipfile.ZipFile(file_path, 'r') as zip_ref:
        for file_info in zip_ref.infolist():
            if not file_info.is_dir():
                file_ext = os.path.splitext(file_info.filename)[1].lower()
                try:
                    with zip_ref.open(file_info.filename) as file:
                        if file_ext == '.pdf':
                            content = load_pdf_content_with_pages(BytesIO(file.read()))
                        elif file_ext == '.pptx':
                            content = load_pptx_content(BytesIO(file.read()))
                        elif file_ext == '.docx':
                            content = load_docx_content(BytesIO(file.read()))
                        elif file_ext in ('.xlsx', '.xls'):
                            content = load_excel_content(BytesIO(file.read()))
                        elif file_ext == '.csv':
                            content = load_csv_content(StringIO(file.read().decode('utf-8')))
                        else:
                            content = None
                            
                        if content is not None:
                            extracted_contents.append({
                                "filename": file_info.filename,
                                "content": content
                            })
                except Exception as e:
                    print(f"Error processing {file_info.filename} in ZIP: {e}")
    return extracted_contents

def load_and_chunk_file(
    file_path: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    strategy: ChunkingStrategy = "recursive",
    **kwargs
) -> List[Dict]:
    """Load and chunk content from various file formats.
    
    Args:
        file_path: Path to the file to be processed.
        chunk_size: Size of each chunk.
        chunk_overlap: Overlap between chunks.
        strategy: Chunking strategy to use.
        **kwargs: Additional arguments for the text splitter.
        
    Returns:
        List of dictionaries containing chunked text and metadata.
    """
    file_content = load_file_content(file_path)
    chunks = []
    
    if isinstance(file_content, list) and all(isinstance(item, dict) and 'page' in item for item in file_content):
        # PDF content with page numbers
        splitter = get_text_splitter(strategy, chunk_size, chunk_overlap, **kwargs)
        for page_content in file_content:
            for chunk in splitter.split_text(page_content['text']):
                chunks.append({
                    "text": chunk,
                    "metadata": {
                        "source": os.path.basename(file_path),
                        "page": page_content['page'],
                        "chunking_strategy": strategy
                    }
                })
    elif isinstance(file_content, list):  # Excel or ZIP with structured data
        for item in file_content:
            if 'data' in item:  # Excel sheet
                for record in item['data']:
                    text = "\n".join(f"{k}: {v}" for k, v in record.items())
                    chunks.append({
                        "text": text,
                        "metadata": {
                            "source": os.path.basename(file_path),
                            "sheet": item.get('sheet_name', ''),
                            "chunking_strategy": "none"  # Excel data isn't chunked
                        }
                    })
            elif 'content' in item:  # ZIP file content
                if isinstance(item['content'], str):
                    # Chunk the text content from files in ZIP
                    splitter = get_text_splitter(strategy, chunk_size, chunk_overlap, **kwargs)
                    for chunk in splitter.split_text(item['content']):
                        chunks.append({
                            "text": chunk,
                            "metadata": {
                                "source": os.path.basename(file_path),
                                "file_in_zip": item['filename'],
                                "chunking_strategy": strategy
                            }
                        })
                elif isinstance(item['content'], list) and all(isinstance(c, dict) and 'page' in c for c in item['content']):
                    # PDF content within ZIP
                    splitter = get_text_splitter(strategy, chunk_size, chunk_overlap, **kwargs)
                    for page_content in item['content']:
                        for chunk in splitter.split_text(page_content['text']):
                            chunks.append({
                                "text": chunk,
                                "metadata": {
                                    "source": os.path.basename(file_path),
                                    "file_in_zip": item['filename'],
                                    "page": page_content['page'],
                                    "chunking_strategy": strategy
                                }
                            })
    else:  # Plain text content from other files
        splitter = get_text_splitter(strategy, chunk_size, chunk_overlap, **kwargs)
        for chunk in splitter.split_text(file_content):
            chunks.append({
                "text": chunk,
                "metadata": {
                    "source": os.path.basename(file_path),
                    "chunking_strategy": strategy
                }
            })
    
    return chunks

def get_text_splitter(
    strategy: ChunkingStrategy,
    chunk_size: int,
    chunk_overlap: int,
    **kwargs
):
    """Get the appropriate text splitter based on strategy."""
    if strategy == "character":
        return CharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separator=kwargs.get("separator", "\n\n")
        )
    elif strategy == "token":
        return TokenTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            encoding_name=kwargs.get("encoding_name", "cl100k_base")
        )
    else:  # recursive
        return RecursiveCharacterTextSplitter(
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
            separators=kwargs.get("separators", ["\n\n", "\n", " ", ""])
        )

# For backward compatibility
def load_and_chunk_pdf(
    file_path: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
    strategy: ChunkingStrategy = "recursive",
    **kwargs
) -> List[Dict]:
    """Legacy function for backward compatibility."""
    return load_and_chunk_file(
        file_path=file_path,
        chunk_size=chunk_size,
        chunk_overlap=chunk_overlap,
        strategy=strategy,
        **kwargs
    )


# from PyPDF2 import PdfReader
# import os
# from langchain.text_splitter import (
#     RecursiveCharacterTextSplitter,
#     CharacterTextSplitter,
#     TokenTextSplitter
# )
# from typing import List, Dict, Literal

# ChunkingStrategy = Literal["recursive", "character", "token"]

# def load_and_chunk_pdf(
#     file_path: str,
#     chunk_size: int = 1000,
#     chunk_overlap: int = 200,
#     strategy: ChunkingStrategy = "recursive",
#     **kwargs
# ) -> List[Dict]:
#     """Handles non-semantic chunking strategies"""
#     reader = PdfReader(file_path)
#     chunks = []
    
#     if strategy == "character":
#         splitter = CharacterTextSplitter(
#             chunk_size=chunk_size,
#             chunk_overlap=chunk_overlap,
#             separator=kwargs.get("separator", "\n\n")
#         )
#     elif strategy == "token":
#         splitter = TokenTextSplitter(
#             chunk_size=chunk_size,
#             chunk_overlap=chunk_overlap,
#             encoding_name=kwargs.get("encoding_name", "cl100k_base")
#         )
#     else:  # recursive
#         splitter = RecursiveCharacterTextSplitter(
#             chunk_size=chunk_size,
#             chunk_overlap=chunk_overlap,
#             separators=kwargs.get("separators", ["\n\n", "\n", " ", ""])
#         )
    
#     for page_num, page in enumerate(reader.pages, start=1):
#         text = page.extract_text()
#         if text.strip():
#             for chunk in splitter.split_text(text):
#                 chunks.append({
#                     "text": chunk,
#                     "metadata": {
#                         "source": os.path.basename(file_path),
#                         "page": page_num,
#                         "chunking_strategy": strategy
#                     }
#                 })
#     return chunks


